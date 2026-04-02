#!/usr/bin/env python3
"""PPT 素材后处理工具

将 SVG 素材转换为高分辨率 PNG，并组装为 PPTX 演示文稿。

默认输出 4K 分辨率（3840×2160, 300DPI），可选 8K 超高分辨率。

使用方式:
    # 完整流水线：SVG → PNG → PPTX（默认 4K）
    python generate_pptx.py ppt_assets/

    # 使用 8K 超高分辨率
    python generate_pptx.py ppt_assets/ --8k

    # 仅转换 PNG
    python generate_pptx.py ppt_assets/ --png-only

    # 仅生成 PPTX（需已有 PNG）
    python generate_pptx.py ppt_assets/ --pptx-only

    # 自定义分辨率和输出文件名
    python generate_pptx.py ppt_assets/ --width 3840 --height 2160 --output my_deck.pptx
"""

import argparse
import glob
import os
import shutil
import subprocess
import sys


def check_rsvg_convert():
    """检查 rsvg-convert 是否可用"""
    path = shutil.which("rsvg-convert")
    if not path:
        print("错误: 未找到 rsvg-convert 工具")
        print("安装方式:")
        print("  macOS:  brew install librsvg")
        print("  Ubuntu: apt install librsvg2-bin")
        sys.exit(1)
    return path


def convert_svg_to_png(
    input_dir: str,
    width: int = 3840,
    height: int = 2160,
    dpi: int = 300,
    background: str = "white",
):
    """将目录下所有 SVG 文件转换为超高分辨率 PNG

    Args:
        input_dir: 包含 SVG 文件的目录
        width: 输出 PNG 宽度（像素）
        height: 输出 PNG 高度（像素）
        dpi: 输出分辨率
        background: 背景颜色
    """
    rsvg = check_rsvg_convert()
    svg_files = sorted(glob.glob(os.path.join(input_dir, "*.svg")))

    if not svg_files:
        print(f"警告: 在 {input_dir} 中未找到 SVG 文件")
        return []

    png_files = []
    for svg_path in svg_files:
        png_path = svg_path.rsplit(".", 1)[0] + ".png"
        svg_name = os.path.basename(svg_path)

        print(f"转换 {svg_name} → 超高分辨率 PNG ({width}x{height}, {dpi}DPI)")

        cmd = [
            rsvg,
            "--format", "png",
            "--width", str(width),
            "--height", str(height),
            "--dpi-x", str(dpi),
            "--dpi-y", str(dpi),
            "--keep-aspect-ratio",
            "--background-color", background,
            svg_path,
            "-o", png_path,
        ]

        try:
            subprocess.run(cmd, check=True, capture_output=True, text=True)
            # 显示文件大小
            size_mb = os.path.getsize(png_path) / (1024 * 1024)
            print(f"  ✓ {os.path.basename(png_path)} ({size_mb:.1f} MB)")
            png_files.append(png_path)
        except subprocess.CalledProcessError as e:
            print(f"  ✗ 转换失败: {e.stderr}")

    print(f"\nPNG 转换完成: {len(png_files)}/{len(svg_files)} 个文件")
    return png_files


def generate_pptx(input_dir: str, output_path: str | None = None):
    """从 PNG 图片生成 PPTX 演示文稿

    每张 PNG 作为一页全出血幻灯片，使用 16:9 宽屏比例。

    Args:
        input_dir: 包含 PNG 文件的目录
        output_path: 输出 PPTX 文件路径，默认为 input_dir/presentation.pptx
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        print("错误: 未安装 python-pptx")
        print("安装: pip install python-pptx")
        sys.exit(1)

    png_files = sorted(glob.glob(os.path.join(input_dir, "*.png")))

    if not png_files:
        print(f"警告: 在 {input_dir} 中未找到 PNG 文件")
        return None

    if output_path is None:
        output_path = os.path.join(input_dir, "presentation.pptx")

    # 创建 16:9 宽屏演示文稿
    prs = Presentation()
    # 16:9 = 13333333 x 7500000 EMU (标准宽屏)
    slide_width = Emu(12192000)   # 25.4cm (标准16:9宽度)
    slide_height = Emu(6858000)   # 14.29cm (标准16:9高度)
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    # 使用空白版式
    blank_layout = prs.slide_layouts[6]

    for png_path in png_files:
        png_name = os.path.basename(png_path)
        print(f"添加幻灯片: {png_name}")

        slide = prs.slides.add_slide(blank_layout)

        # 图片全出血覆盖幻灯片
        slide.shapes.add_picture(
            png_path,
            left=Emu(0),
            top=Emu(0),
            width=slide_width,
            height=slide_height,
        )

    prs.save(output_path)
    size_mb = os.path.getsize(output_path) / (1024 * 1024)
    print(f"\nPPTX 生成完成: {output_path} ({size_mb:.1f} MB, {len(png_files)} 页)")
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="PPT 素材后处理工具 - SVG → PNG → PPTX",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  # 完整流水线（默认 4K）
  python generate_pptx.py ppt_assets/

  # 使用 8K 超高分辨率
  python generate_pptx.py ppt_assets/ --8k

  # 仅转换 PNG
  python generate_pptx.py ppt_assets/ --png-only

  # 仅生成 PPTX
  python generate_pptx.py ppt_assets/ --pptx-only

  # 自定义分辨率
  python generate_pptx.py ppt_assets/ --width 3840 --height 2160 --dpi 300
        """,
    )

    parser.add_argument(
        "input_dir",
        help="包含 SVG/PNG 文件的目录（默认: ppt_assets/）",
        nargs="?",
        default="ppt_assets",
    )
    parser.add_argument(
        "--output", "-o",
        help="输出 PPTX 文件路径（默认: <input_dir>/presentation.pptx）",
    )
    parser.add_argument(
        "--png-only",
        action="store_true",
        help="仅执行 SVG → PNG 转换",
    )
    parser.add_argument(
        "--pptx-only",
        action="store_true",
        help="仅执行 PNG → PPTX 生成（需已有 PNG）",
    )
    parser.add_argument(
        "--8k",
        action="store_true",
        dest="ultra_hd",
        help="使用 8K 超高分辨率（7680×4320, 600DPI），覆盖 --width/--height/--dpi",
    )
    parser.add_argument(
        "--width", "-W",
        type=int,
        default=3840,
        help="PNG 输出宽度，像素（默认: 3840，4K）",
    )
    parser.add_argument(
        "--height", "-H",
        type=int,
        default=2160,
        help="PNG 输出高度，像素（默认: 2160，4K）",
    )
    parser.add_argument(
        "--dpi",
        type=int,
        default=300,
        help="PNG 输出 DPI（默认: 300）",
    )
    parser.add_argument(
        "--background", "-bg",
        default="white",
        help="PNG 背景颜色（默认: white）",
    )

    args = parser.parse_args()

    # --8k 快捷选项覆盖分辨率参数
    if args.ultra_hd:
        args.width = 7680
        args.height = 4320
        args.dpi = 600

    if not os.path.isdir(args.input_dir):
        print(f"错误: 目录不存在 - {args.input_dir}")
        sys.exit(1)

    print(f"=== PPT 素材后处理工具 ===")
    print(f"工作目录: {os.path.abspath(args.input_dir)}\n")

    if args.pptx_only:
        # 仅生成 PPTX
        generate_pptx(args.input_dir, args.output)
    elif args.png_only:
        # 仅转换 PNG
        convert_svg_to_png(
            args.input_dir,
            width=args.width,
            height=args.height,
            dpi=args.dpi,
            background=args.background,
        )
    else:
        # 完整流水线：SVG → PNG → PPTX
        png_files = convert_svg_to_png(
            args.input_dir,
            width=args.width,
            height=args.height,
            dpi=args.dpi,
            background=args.background,
        )
        if png_files:
            print()  # 空行分隔
            generate_pptx(args.input_dir, args.output)

    print("\n完成!")


if __name__ == "__main__":
    main()
